"""PPTX comments fixture generator.

``python-pptx`` doesn't author comment parts. The fixture is produced by
building a vanilla 3-slide deck with ``python-pptx`` and then injecting
``ppt/commentAuthors.xml`` plus ``ppt/comments/comment{N}.xml`` parts
into the resulting zip (one comments file per slide that has comments).

Shape mirrors what ``crates/xberg/src/extraction/pptx/comments.rs``
consumes: ``<p:cmAuthor id="…" name="…"/>`` for authors,
``<p:cm authorId="…" dt="…" idx="…"><p:text>…</p:text></p:cm>`` for
comments. The extractor anchors the resulting ``DocumentRevision`` on
``RevisionAnchor::Slide { index }`` where index is the zero-based slide
ordinal (so ``comment1.xml`` -> slide index 0, ``comment3.xml`` -> 2).
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from pptx import Presentation  # type: ignore[import-untyped, import-not-found, unused-ignore]

from .gt_schema import revisions_expectation, write_ground_truth

ZIP_MTIME = (2024, 1, 1, 0, 0, 0)

AUTHORS = [
    (0, "Alice"),
    (1, "Bob"),
]

# Each row: (slide_index_zero_based, idx, author_id, dt, text)
COMMENTS = [
    (0, 1, 0, "2024-06-01T10:00:00Z", "Alice: opening question on slide 1"),
    (0, 2, 1, "2024-06-01T10:15:00Z", "Bob: follow-up on slide 1"),
    (2, 1, 0, "2024-06-01T11:30:00Z", "Alice: closing comment on slide 3"),
]


def _build_baseline_pptx() -> bytes:
    """Author a vanilla 3-slide deck with one text shape each."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    for i in range(3):
        slide = prs.slides.add_slide(blank_layout)
        # python-pptx writes deterministic slideN.xml; add a minimal text
        # frame so each slide carries body text.
        textbox = slide.shapes.add_textbox(left=914400, top=914400, width=914400 * 4, height=914400)
        textbox.text_frame.text = f"Slide {i + 1} body"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _comment_authors_xml() -> bytes:
    """Build ``ppt/commentAuthors.xml``."""
    authors_xml = "".join(
        f'<p:cmAuthor id="{aid}" name="{name}" initials="{name[0]}" lastIdx="0" clrIdx="0"/>' for aid, name in AUTHORS
    )
    xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        f"{authors_xml}"
        "</p:cmAuthorLst>"
    )
    return xml.encode("utf-8")


def _comments_for_slide(slide_index: int) -> bytes | None:
    """Return ``ppt/comments/comment{slide_index+1}.xml`` bytes, or ``None``."""
    entries = [c for c in COMMENTS if c[0] == slide_index]
    if not entries:
        return None
    inner = "".join(
        f'<p:cm authorId="{aid}" dt="{dt}" idx="{idx}"><p:text>{text}</p:text></p:cm>'
        for (_, idx, aid, dt, text) in entries
    )
    xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        f"{inner}"
        "</p:cmLst>"
    )
    return xml.encode("utf-8")


def _patch_content_types(original: bytes, comment_slide_indices: list[int]) -> bytes:
    """Register commentAuthors + per-slide comments content-types."""
    text = original.decode("utf-8")
    overrides: list[str] = [
        '<Override PartName="/ppt/commentAuthors.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"/>'
    ]
    for slide_idx in comment_slide_indices:
        overrides.append(
            f'<Override PartName="/ppt/comments/comment{slide_idx + 1}.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.presentationml.comments+xml"/>'
        )
    addition = "".join(o for o in overrides if o not in text)
    if not addition:
        return original
    return text.replace("</Types>", f"{addition}</Types>").encode("utf-8")


def _rewrite_zip(src_bytes: bytes, additions: dict[str, bytes], replacements: dict[str, bytes]) -> bytes:
    """Re-zip with deterministic mtimes; additions are appended after the original entries."""
    buf = io.BytesIO()
    seen: set[str] = set()
    with zipfile.ZipFile(io.BytesIO(src_bytes), "r") as src:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as dst:
            for name in src.namelist():
                data = replacements.get(name, src.read(name))
                info = zipfile.ZipInfo(name, ZIP_MTIME)
                info.compress_type = zipfile.ZIP_DEFLATED
                dst.writestr(info, data)
                seen.add(name)
            for name, data in additions.items():
                if name in seen:
                    continue
                info = zipfile.ZipInfo(name, ZIP_MTIME)
                info.compress_type = zipfile.ZIP_DEFLATED
                dst.writestr(info, data)
    return buf.getvalue()


def generate(output_root: Path, repo_root: Path) -> list[Path]:
    """Emit pptx_comments_basic.pptx + sidecar under ``output_root/pptx/``."""
    output_dir = output_root / "pptx"
    output_dir.mkdir(parents=True, exist_ok=True)

    base = _build_baseline_pptx()
    comment_slide_indices = sorted({c[0] for c in COMMENTS})

    additions: dict[str, bytes] = {"ppt/commentAuthors.xml": _comment_authors_xml()}
    for slide_idx in comment_slide_indices:
        payload = _comments_for_slide(slide_idx)
        assert payload is not None  # by construction
        additions[f"ppt/comments/comment{slide_idx + 1}.xml"] = payload

    with zipfile.ZipFile(io.BytesIO(base), "r") as zf:
        content_types = zf.read("[Content_Types].xml")
    replacements = {"[Content_Types].xml": _patch_content_types(content_types, comment_slide_indices)}

    out = _rewrite_zip(base, additions=additions, replacements=replacements)

    fixture_path = output_dir / "pptx_comments_basic.pptx"
    sidecar_path = output_dir / "pptx_comments_basic.gt.json"
    fixture_path.write_bytes(out)

    expected_revisions = [
        {
            "kind": "Comment",
            "author": dict(AUTHORS)[author_id],
            "timestamp": dt,
            "slide_index": slide_idx,
        }
        for (slide_idx, _idx, author_id, dt, _text) in COMMENTS
    ]

    write_ground_truth(
        sidecar_path,
        fixture_path,
        repo_root,
        document_format="pptx",
        feature="revisions",
        expectations=revisions_expectation(
            expected_count=len(COMMENTS),
            revisions=expected_revisions,
            notes=(
                "Three slides; comments on slide 1 (two by Alice and Bob) and slide 3 "
                "(one by Alice). RevisionKind = Comment for every entry; anchor is "
                "RevisionAnchor::Slide with zero-based index. Note that this fixture "
                "intentionally does not include the slide -> comments .rels link — the "
                "extractor walks the comments directory by filename pattern, not via the "
                "relationship graph."
            ),
        ),
        generator="pptx_comments",
    )
    return [fixture_path, sidecar_path]
